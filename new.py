import streamlit as st

# --- MUST BE THE FIRST STREAMLIT COMMAND ---
st.set_page_config(page_title="Hybrid AI Presentation Architect", page_icon="📊", layout="wide")

import requests
import json
import io
import re
import os
from pptx import Presentation
from pptx.util import Pt
from pydantic import BaseModel, Field
from typing import List
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

# Import Agno for Gemini
try:
    from agno.agent import Agent
    from agno.models.google import Gemini
except ImportError:
    st.sidebar.error("⚠️ Library Missing: Run 'pip install agno google-generativeai'")

# --- App Constants ---
CHAT_API_URL = "https://api.sarvam.ai/v1/chat/completions"
TRANSLATE_API_URL = "https://api.sarvam.ai/translate"

SUPPORTED_LANGUAGES = {
    'en-IN': 'English',
    'hi-IN': 'Hindi',
    'ta-IN': 'Tamil',
    'te-IN': 'Telugu',
    'bn-IN': 'Bengali',
    'kn-IN': 'Kannada',
    'mr-IN': 'Marathi',
    'gu-IN': 'Gujarati'
}

# --- Pydantic Models for Agno/Gemini ---
class SlideContentItem(BaseModel):
    type: str = Field(description="Must be 'subheading', 'point', or 'paragraph'")
    text: str = Field(description="The text content")

class SlideItem(BaseModel):
    heading: str = Field(description="The main heading of the slide")
    content: List[SlideContentItem] = Field(description="List of items for this slide")

class PresentationData(BaseModel):
    slides: List[SlideItem] = Field(description="List of all slides to be generated")

# --- Formatting & Cleanup Helper Functions ---
def sanitize_text(text):
    """Removes invisible control characters that can corrupt PowerPoint XML."""
    if not isinstance(text, str):
        return str(text)
    return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)

def apply_formatting(run, font_name='Times New Roman', size=12, bold=False):
    """Applies specific font, size, and boldness to a python-pptx Run."""
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold

def clean_and_parse_json(raw: str):
    """Robustly clean AI-generated text and parse it as JSON."""
    # 1. Strip markdown code fences
    txt = re.sub(r'```json|```', '', raw).strip()
    # 2. Replace smart/curly quotes with straight quotes
    txt = txt.replace('\u201c', '"').replace('\u201d', '"').replace('\u2018', "'").replace('\u2019', "'")
    # 3. Collapse mid-word newlines and strip tabs/carriage returns
    txt = re.sub(r'(?<=\S)\n(?=\S)', ' ', txt)
    txt = re.sub(r'[\t\r]', ' ', txt)
    # 4. Remove trailing commas before ] or } (invalid JSON)
    txt = re.sub(r',\s*([\]}])', r'\1', txt)
    # 5. Extract the outermost JSON object or array
    match = re.search(r'(\[.*\]|\{.*\})', txt, re.DOTALL)
    if not match:
        raise ValueError("No JSON structure found in AI response.")
    return json.loads(match.group(0))

# --- Sarvam API helper (single call) ---
def _sarvam_chat(prompt: str, api_key: str) -> str:
    """Makes a single call to the Sarvam chat API and returns the raw text response."""
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {"model": "sarvam-m", "messages": [{"role": "user", "content": prompt}]}
    response = requests.post(CHAT_API_URL, headers=headers, json=payload)
    response.raise_for_status()
    return response.json()["choices"][0]["message"]["content"]

# --- Backend AI Functions ---

def generate_gemini_presentation(doc_text: str, slide_instructions: str, api_key: str, slide_count: int, content_length: str) -> list:
    """Generates structured presentation content using Google Gemini via Agno."""
    os.environ["GOOGLE_API_KEY"] = api_key

    length_instruction = {
        "Brief": "Keep content short. Generate 4-6 points PER SLIDE.",
        "Medium": "Moderate detail. Generate 6-10 points PER SLIDE.",
        "Detailed": "Comprehensive detail. Generate 12-15 points PER SLIDE."
    }[content_length]

    instructions = [
        f"Create a professional {slide_count}-slide presentation based on the provided documentation.",
        f"IMPORTANT: You MUST generate exactly {slide_count} individual slides.",
        f"Content Density: {length_instruction}",
        "CRITICAL JSON RULES:",
        "1. Output must be a valid JSON object matching the PresentationData schema.",
        "2. No conversational text or preamble. No markdown code blocks like ```json.",
        "3. Use single quotes for inner text if needed to avoid breaking the JSON structure."
    ]

    if slide_instructions:
        instructions.append(f"Follow these slide-specific instructions: {slide_instructions}")

    agent = Agent(
        model=Gemini(id="gemini-2.0-flash-lite"),
        description="Expert presentation content creator.",
        instructions=instructions
    )

    response = agent.run(f"Generate the full {slide_count}-slide presentation structure using this documentation: {doc_text}")
    presentation_data = response.content

    # Robust Fallback Parsing
    if isinstance(presentation_data, str):
        try:
            data = clean_and_parse_json(presentation_data)
            if isinstance(data, list):
                presentation_data = PresentationData(slides=data)
            else:
                presentation_data = PresentationData(**data)
        except Exception as e:
            st.error(f"DEBUG: AI returned raw text that failed parsing: {str(presentation_data)[:300]}")
            raise ValueError(f"AI returned invalid data format: {str(e)}")

    if not hasattr(presentation_data, 'slides'):
        raise ValueError("AI response is missing the 'slides' list.")

    return [
        {"heading": s.heading, "content": [{"type": i.type, "text": i.text} for i in s.content]}
        for s in presentation_data.slides
    ]


def generate_sarvam_presentation(doc_text: str, slide_instructions: str, api_key: str, slide_count: int, content_length: str) -> list:
    """
    Multi-agent Sarvam pipeline:
      Stage 1 (Planner): Generate slide titles/outline.
      Stage 2 (Writers): One sub-agent per slide generates only that slide's content.
    """
    length_map = {
        "Brief": "4-6 bullet points",
        "Medium": "6-10 bullet points",
        "Detailed": "12-15 bullet points"
    }
    length_instruction = length_map[content_length]

    # ── Stage 1: Planner ─────────────────────────────────────────────────────
    # Parse any user slide instructions into a list keyed by slide number
    user_slide_hints = {}
    if slide_instructions:
        for line in slide_instructions.strip().splitlines():
            m = re.match(r'[Ss]lide\s*(\d+)[:\-–]\s*(.*)', line)
            if m:
                user_slide_hints[int(m.group(1))] = m.group(2).strip()

    planner_prompt = f"""
You are a presentation planner. Based on the documentation below, generate exactly {slide_count} slide titles.
{"The user has provided hints: " + slide_instructions if slide_instructions else ""}

Documentation:
{doc_text[:3000]}

Output ONLY a valid JSON array of {slide_count} strings representing slide titles.
Example: ["Introduction", "Key Features", "Conclusion"]
No preamble, no extra text.
"""
    planner_raw = _sarvam_chat(planner_prompt, api_key)
    try:
        titles = clean_and_parse_json(planner_raw)
        if not isinstance(titles, list):
            raise ValueError("Planner did not return a list.")
    except Exception:
        # Fallback: extract quoted strings
        titles = re.findall(r'"([^"]+)"', planner_raw)

    # Pad or trim to exact count
    while len(titles) < slide_count:
        titles.append(f"Slide {len(titles) + 1}")
    titles = titles[:slide_count]

    # ── Stage 2: Per-Slide Writers ────────────────────────────────────────────
    slides = []
    progress = st.progress(0, text="Writing slides...")

    for i, title in enumerate(titles):
        hint = user_slide_hints.get(i + 1, "")
        hint_text = f'Focus: "{hint}"' if hint else ""

        writer_prompt = f"""
You are a slide content writer. Write content for slide {i + 1} of {slide_count}.
Slide title: "{title}"
{hint_text}

Content from the documentation to draw from:
{doc_text[:3000]}

Write {length_instruction} for this slide.
Output ONLY a valid JSON array of content items.
Each item must follow this schema: {{"type": "point", "text": "..."}}
Use "type": "subheading" for sub-sections, "type": "paragraph" for prose.
No preamble, no extra text. Output the JSON array only.
"""
        writer_raw = _sarvam_chat(writer_prompt, api_key)
        try:
            content_items = clean_and_parse_json(writer_raw)
            if not isinstance(content_items, list):
                content_items = [{"type": "point", "text": str(content_items)}]
        except Exception:
            content_items = [{"type": "point", "text": writer_raw.strip()[:300]}]

        slides.append({"heading": title, "content": content_items})
        progress.progress((i + 1) / slide_count, text=f"Wrote slide {i + 1}/{slide_count}: {title}")

    progress.empty()
    return slides


def translate_content(text: str, target_lang: str, api_key: str) -> str:
    """Translates text to the target language using Sarvam."""
    if not text.strip() or target_lang == 'en-IN':
        return text
    headers = {"api-subscription-key": api_key, "Content-Type": "application/json"}
    payload = {"input": text, "source_language_code": "en-IN", "target_language_code": target_lang}
    res = requests.post(TRANSLATE_API_URL, headers=headers, json=payload)
    res.raise_for_status()
    return res.json()["translated_text"]


# --- UI Logic ---
st.title("🤖 Hybrid AI Presentation Architect")
st.markdown("Build professional presentations using **Sarvam AI** or **Google Gemini**.")

with st.sidebar:
    st.header("⚙️ Configuration")
    ai_model = st.selectbox("AI Model", options=["Google Gemini", "Sarvam AI"])
    language = st.selectbox("Language", options=list(SUPPORTED_LANGUAGES.keys()), format_func=lambda x: SUPPORTED_LANGUAGES[x])
    slide_count = st.slider("Slides", 2, 10, 5)

    st.header("🔑 API Keys")
    sarvam_key = st.text_input("Sarvam Key", type="password") if (ai_model == "Sarvam AI" or language != 'en-IN') else ""
    gemini_key = st.text_input("Gemini Key", type="password") if ai_model == "Google Gemini" else ""

    st.header("📝 Style")
    content_length = st.selectbox("Length", options=["Brief", "Medium", "Detailed"], index=1)

# Main Form
col1, col2 = st.columns(2)
with col1:
    uploaded_doc = st.file_uploader("📂 Upload Documentation (.txt, .md, .docx)", type=["txt", "md", "docx"])
with col2:
    uploaded_template = st.file_uploader("🎨 Upload PPT Template (.pptx)", type=["pptx"])

slide_instructions = st.text_area(
    "📝 Slide-by-Slide Instructions (Optional)",
    placeholder="Slide 1: Introduction and overview\nSlide 2: Technical architecture details\nSlide 3: Benefits and use cases",
    height=150
)

if not uploaded_doc or not uploaded_template:
    st.info("👋 Welcome! Please upload your **Documentation** and **Sample PPTX** to get started.")
    st.stop()

if st.button("🚀 Build My Presentation", type="primary"):
    if ai_model == "Google Gemini" and not gemini_key:
        st.error("Missing Gemini API Key.")
    elif (ai_model == "Sarvam AI" or language != 'en-IN') and not sarvam_key:
        st.error("Missing Sarvam API Key.")
    else:
        try:
            # Read documentation content
            file_name = uploaded_doc.name.lower()
            if file_name.endswith(".docx"):
                if DocxDocument is None:
                    st.error("python-docx is not installed. Run: pip install python-docx")
                    st.stop()
                docx_obj = DocxDocument(io.BytesIO(uploaded_doc.read()))
                paragraphs = [p.text for p in docx_obj.paragraphs if p.text.strip()]
                # Also extract text from tables
                for table in docx_obj.tables:
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            paragraphs.append(row_text)
                doc_content = "\n".join(paragraphs)
            else:
                doc_content = uploaded_doc.read().decode("utf-8")

            if ai_model == "Sarvam AI":
                st.info(f"🤖 Using multi-agent pipeline: 1 planner + {slide_count} slide writers")
                final_slides = generate_sarvam_presentation(doc_content, slide_instructions, sarvam_key, slide_count, content_length)
            else:
                with st.spinner(f"🧠 Generating {slide_count} slides via Google Gemini..."):
                    final_slides = generate_gemini_presentation(doc_content, slide_instructions, gemini_key, slide_count, content_length)

            # Ensure we didn't get fewer slides than requested
            if len(final_slides) < slide_count:
                st.warning(f"Note: AI generated {len(final_slides)} slides instead of {slide_count}. Processing anyway...")

            if language != 'en-IN':
                with st.spinner(f"🌐 Translating to {SUPPORTED_LANGUAGES[language]}..."):
                    for s in final_slides:
                        s['heading'] = translate_content(s['heading'], language, sarvam_key)
                        for item in s['content']:
                            item['text'] = translate_content(item['text'], language, sarvam_key)

            with st.spinner("🎨 Applying styles to template..."):
                prs = Presentation(io.BytesIO(uploaded_template.read()))

                # Logic: Reuse template slides 1..N (preserving slide 0)
                for i, s_data in enumerate(final_slides):
                    idx = i + 1
                    if idx < len(prs.slides):
                        slide = prs.slides[idx]
                    else:
                        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
                        slide = prs.slides.add_slide(layout)

                    if slide.shapes.title:
                        slide.shapes.title.text = sanitize_text(s_data['heading'])
                        for p in slide.shapes.title.text_frame.paragraphs:
                            for r in p.runs:
                                apply_formatting(r, size=18, bold=True)

                    body = next((sh for sh in slide.shapes.placeholders if sh.placeholder_format.idx == 1), None)
                    if body and body.has_text_frame:
                        tf = body.text_frame
                        tf.clear()
                        for j, item in enumerate(s_data['content']):
                            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                            p.text = sanitize_text(item['text'])

                            # Layout logic
                            is_para_or_sub = item['type'] in ['subheading', 'paragraph']
                            p.level = 0 if is_para_or_sub else 1

                            if p.runs:
                                size = 14 if item['type'] == 'subheading' else 12
                                bold = item['type'] == 'subheading'
                                apply_formatting(p.runs[0], size=size, bold=bold)

                # Cleanup extra original slides
                while len(prs.slides) > (len(final_slides) + 1):
                    rId = prs.slides._sldIdLst[-1].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[-1]

                out = io.BytesIO()
                prs.save(out)
                out.seek(0)

            st.success("✅ Presentation ready!")
            st.download_button("📥 Download PPTX", data=out, file_name="presentation.pptx")
            st.balloons()

        except Exception as e:
            st.error(f"Error: {e}")