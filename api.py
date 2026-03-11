"""
FastAPI Backend for Hybrid AI Presentation Architect
=====================================================
Exposes the same AI logic from new.py as a REST API.

Endpoints:
  GET  /                  - Health check
  GET  /languages         - List supported languages
  POST /generate/slides   - Generate slide JSON from documentation
  POST /generate/pptx     - Generate and download a .pptx file

Run with:
  uvicorn api:app --reload
  
Docs at:
  http://localhost:8000/docs
"""

import json
import io
import re
import os
from typing import Optional, List

import requests
from fastapi import FastAPI, File, Form, UploadFile, HTTPException
from fastapi.responses import StreamingResponse, JSONResponse
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Pt

# --- Agno / Gemini (optional) ---
try:
    from agno.agent import Agent
    from agno.models.google import Gemini
    AGNO_AVAILABLE = True
except ImportError:
    AGNO_AVAILABLE = False

# ─────────────────────────────────────────────────────────────────────────────
# Constants
# ─────────────────────────────────────────────────────────────────────────────

CHAT_API_URL = "https://api.sarvam.ai/v1/chat/completions"
TRANSLATE_API_URL = "https://api.sarvam.ai/translate"

SUPPORTED_LANGUAGES = {
    "en-IN": "English",
    "hi-IN": "Hindi",
    "ta-IN": "Tamil",
    "te-IN": "Telugu",
    "bn-IN": "Bengali",
    "kn-IN": "Kannada",
    "mr-IN": "Marathi",
    "gu-IN": "Gujarati",
}

CONTENT_LENGTH_MAP = {
    "Brief": {
        "gemini": "Keep content short. Generate 4-6 points PER SLIDE.",
        "sarvam": "4-6 bullet points",
    },
    "Medium": {
        "gemini": "Moderate detail. Generate 6-10 points PER SLIDE.",
        "sarvam": "6-10 bullet points",
    },
    "Detailed": {
        "gemini": "Comprehensive detail. Generate 12-15 points PER SLIDE.",
        "sarvam": "12-15 bullet points",
    },
}

# ─────────────────────────────────────────────────────────────────────────────
# Pydantic Response Models
# ─────────────────────────────────────────────────────────────────────────────

class ContentItem(BaseModel):
    type: str = Field(description="'point', 'subheading', or 'paragraph'")
    text: str

class Slide(BaseModel):
    heading: str
    content: List[ContentItem]

class SlidesResponse(BaseModel):
    slides: List[Slide]
    slide_count: int
    ai_model: str

# Pydantic models used for Gemini structured output
class SlideContentItem(BaseModel):
    type: str = Field(description="Must be 'subheading', 'point', or 'paragraph'")
    text: str = Field(description="The text content")

class SlideItem(BaseModel):
    heading: str = Field(description="The main heading of the slide")
    content: List[SlideContentItem] = Field(description="List of items for this slide")

class PresentationData(BaseModel):
    slides: List[SlideItem] = Field(description="List of all slides to be generated")


# ─────────────────────────────────────────────────────────────────────────────
# Helper Utilities (mirror of new.py)
# ─────────────────────────────────────────────────────────────────────────────

def sanitize_text(text: str) -> str:
    """Removes invisible control characters that can corrupt PowerPoint XML."""
    if not isinstance(text, str):
        return str(text)
    return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)


def apply_formatting(run, font_name: str = 'Times New Roman', size: int = 12, bold: bool = False):
    """Applies specific font, size, and boldness to a python-pptx Run."""
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold


def clean_and_parse_json(raw: str):
    """Robustly clean AI-generated text and parse it as JSON."""
    txt = re.sub(r'```json|```', '', raw).strip()
    txt = txt.replace('\u201c', '"').replace('\u201d', '"').replace('\u2018', "'").replace('\u2019', "'")
    txt = re.sub(r'(?<=\S)\n(?=\S)', ' ', txt)
    txt = re.sub(r'[\t\r]', ' ', txt)
    txt = re.sub(r',\s*([\]}])', r'\1', txt)
    match = re.search(r'(\[.*\]|\{.*\})', txt, re.DOTALL)
    if not match:
        raise ValueError("No JSON structure found in AI response.")
    return json.loads(match.group(0))


def _sarvam_chat(prompt: str, api_key: str) -> str:
    """Makes a single call to the Sarvam chat API."""
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {"model": "sarvam-m", "messages": [{"role": "user", "content": prompt}]}
    response = requests.post(CHAT_API_URL, headers=headers, json=payload, timeout=60)
    response.raise_for_status()
    return response.json()["choices"][0]["message"]["content"]


def translate_content(text: str, target_lang: str, api_key: str) -> str:
    """Translates text to the target language using Sarvam Translate API."""
    if not text.strip() or target_lang == 'en-IN':
        return text
    headers = {"api-subscription-key": api_key, "Content-Type": "application/json"}
    payload = {"input": text, "source_language_code": "en-IN", "target_language_code": target_lang}
    res = requests.post(TRANSLATE_API_URL, headers=headers, json=payload, timeout=30)
    res.raise_for_status()
    return res.json()["translated_text"]


# ─────────────────────────────────────────────────────────────────────────────
# AI Generation Logic
# ─────────────────────────────────────────────────────────────────────────────

def _generate_gemini(doc_text: str, slide_instructions: str, api_key: str, slide_count: int, content_length: str) -> list:
    """Generates slides using Google Gemini via Agno."""
    if not AGNO_AVAILABLE:
        raise HTTPException(status_code=500, detail="Agno library not installed. Run: pip install agno google-generativeai")

    os.environ["GOOGLE_API_KEY"] = api_key
    length_instruction = CONTENT_LENGTH_MAP[content_length]["gemini"]

    instructions = [
        f"Create a professional {slide_count}-slide presentation based on the provided documentation.",
        f"IMPORTANT: You MUST generate exactly {slide_count} individual slides.",
        f"Content Density: {length_instruction}",
        "CRITICAL JSON RULES:",
        "1. Output must be a valid JSON object matching the PresentationData schema.",
        "2. No conversational text or preamble. No markdown code blocks like ```json.",
        "3. Use single quotes for inner text if needed to avoid breaking the JSON structure.",
    ]
    if slide_instructions:
        instructions.append(f"Follow these slide-specific instructions: {slide_instructions}")

    agent = Agent(
        model=Gemini(id="gemini-2.0-flash-lite"),
        description="Expert presentation content creator.",
        instructions=instructions,
    )

    response = agent.run(f"Generate the full {slide_count}-slide presentation using this documentation: {doc_text}")
    presentation_data = response.content

    if isinstance(presentation_data, str):
        data = clean_and_parse_json(presentation_data)
        if isinstance(data, list):
            presentation_data = PresentationData(slides=data)
        else:
            presentation_data = PresentationData(**data)

    if not hasattr(presentation_data, 'slides'):
        raise ValueError("AI response is missing the 'slides' list.")

    return [
        {"heading": s.heading, "content": [{"type": i.type, "text": i.text} for i in s.content]}
        for s in presentation_data.slides
    ]


def _generate_sarvam(doc_text: str, slide_instructions: str, api_key: str, slide_count: int, content_length: str) -> list:
    """
    Multi-agent Sarvam pipeline:
      Stage 1 (Planner): Generate slide titles.
      Stage 2 (Writers): One sub-agent per slide writes the content.
    """
    length_instruction = CONTENT_LENGTH_MAP[content_length]["sarvam"]

    # Parse per-slide instructions
    user_slide_hints = {}
    if slide_instructions:
        for line in slide_instructions.strip().splitlines():
            m = re.match(r'[Ss]lide\s*(\d+)[:\-–]\s*(.*)', line)
            if m:
                user_slide_hints[int(m.group(1))] = m.group(2).strip()

    # ── Stage 1: Planner ──────────────────────────────────────────────────────
    planner_prompt = f"""
You are a presentation planner. Based on the documentation below, generate exactly {slide_count} slide titles.
{"User hints: " + slide_instructions if slide_instructions else ""}

Documentation:
{doc_text[:3000]}

Output ONLY a valid JSON array of {slide_count} strings representing slide titles.
Example: ["Introduction", "Key Features", "Conclusion"]
No preamble, no extra text.
"""
    planner_raw = _sarvam_chat(planner_prompt, api_key)
    try:
        titles = clean_and_parse_json(planner_raw)
        if not isinstance(titles, list):
            raise ValueError("Planner did not return a list.")
    except Exception:
        titles = re.findall(r'"([^"]+)"', planner_raw)

    while len(titles) < slide_count:
        titles.append(f"Slide {len(titles) + 1}")
    titles = titles[:slide_count]

    # ── Stage 2: Per-Slide Writers ────────────────────────────────────────────
    slides = []
    for i, title in enumerate(titles):
        hint = user_slide_hints.get(i + 1, "")
        hint_text = f'Focus: "{hint}"' if hint else ""

        writer_prompt = f"""
You are a slide content writer. Write content for slide {i + 1} of {slide_count}.
Slide title: "{title}"
{hint_text}

Documentation:
{doc_text[:3000]}

Write {length_instruction} for this slide.
Output ONLY a valid JSON array of content items.
Schema: [{{"type": "point", "text": "..."}}]
Use "type": "subheading" for sub-sections, "type": "paragraph" for prose.
No preamble, no extra text.
"""
        writer_raw = _sarvam_chat(writer_prompt, api_key)
        try:
            content_items = clean_and_parse_json(writer_raw)
            if not isinstance(content_items, list):
                content_items = [{"type": "point", "text": str(content_items)}]
        except Exception:
            content_items = [{"type": "point", "text": writer_raw.strip()[:300]}]

        slides.append({"heading": title, "content": content_items})

    return slides


def _build_pptx(slides: list, template_bytes: bytes) -> bytes:
    """Fills a PPTX template with slide content and returns raw bytes."""
    prs = Presentation(io.BytesIO(template_bytes))

    for i, s_data in enumerate(slides):
        idx = i + 1
        if idx < len(prs.slides):
            slide = prs.slides[idx]
        else:
            layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = sanitize_text(s_data['heading'])
            for p in slide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    apply_formatting(r, size=18, bold=True)

        body = next((sh for sh in slide.shapes.placeholders if sh.placeholder_format.idx == 1), None)
        if body and body.has_text_frame:
            tf = body.text_frame
            tf.clear()
            for j, item in enumerate(s_data['content']):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = sanitize_text(item['text'])
                is_para_or_sub = item['type'] in ['subheading', 'paragraph']
                p.level = 0 if is_para_or_sub else 1
                if p.runs:
                    size = 14 if item['type'] == 'subheading' else 12
                    bold = item['type'] == 'subheading'
                    apply_formatting(p.runs[0], size=size, bold=bold)

    # Remove extra template slides beyond slide 0 + generated slides
    while len(prs.slides) > (len(slides) + 1):
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[-1]

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()


def _run_generation(
    doc_text: str,
    slide_instructions: str,
    ai_model: str,
    slide_count: int,
    content_length: str,
    sarvam_key: str,
    gemini_key: str,
) -> list:
    """Dispatches to the correct AI backend and returns a list of slide dicts."""
    if ai_model == "Sarvam AI":
        if not sarvam_key:
            raise HTTPException(status_code=422, detail="sarvam_key is required for Sarvam AI.")
        return _generate_sarvam(doc_text, slide_instructions, sarvam_key, slide_count, content_length)
    elif ai_model == "Google Gemini":
        if not gemini_key:
            raise HTTPException(status_code=422, detail="gemini_key is required for Google Gemini.")
        return _generate_gemini(doc_text, slide_instructions, gemini_key, slide_count, content_length)
    else:
        raise HTTPException(status_code=422, detail=f"Unknown ai_model: '{ai_model}'. Use 'Sarvam AI' or 'Google Gemini'.")


# ─────────────────────────────────────────────────────────────────────────────
# FastAPI App
# ─────────────────────────────────────────────────────────────────────────────

app = FastAPI(
    title="Hybrid AI Presentation Architect",
    description="Generates professional .pptx presentations from documentation using Sarvam AI or Google Gemini.",
    version="1.0.0",
)


@app.get("/", tags=["Health"])
def root():
    return {"status": "ok", "message": "Presentation API is running. Visit /docs for the full API."}


@app.get("/languages", tags=["Info"])
def get_languages():
    """Returns all supported language codes and names for translation."""
    return SUPPORTED_LANGUAGES


@app.post("/generate/slides", response_model=SlidesResponse, tags=["Generate"])
async def generate_slides(
    doc_file: UploadFile = File(..., description="Documentation file (.txt or .md)"),
    slide_count: int = Form(5, ge=2, le=10, description="Number of slides to generate"),
    content_length: str = Form("Medium", description="'Brief', 'Medium', or 'Detailed'"),
    ai_model: str = Form("Sarvam AI", description="'Sarvam AI' or 'Google Gemini'"),
    sarvam_key: str = Form("", description="Sarvam API key (required for Sarvam AI or translation)"),
    gemini_key: str = Form("", description="Google Gemini API key (required for Google Gemini)"),
    slide_instructions: str = Form("", description="Optional per-slide instructions"),
    target_language: str = Form("en-IN", description="Target language code (e.g. 'hi-IN')"),
):
    """
    Generate presentation slide content as JSON from uploaded documentation.
    
    - Upload a `.txt` or `.md` file as the documentation source.
    - Optionally provide slide-by-slide instructions.
    - Returns a JSON list of slides with heading and content items.
    """
    if content_length not in CONTENT_LENGTH_MAP:
        raise HTTPException(status_code=422, detail=f"content_length must be one of: {list(CONTENT_LENGTH_MAP.keys())}")
    if target_language not in SUPPORTED_LANGUAGES:
        raise HTTPException(status_code=422, detail=f"Unsupported language code. Use /languages to see valid options.")

    raw_bytes = await doc_file.read()
    try:
        doc_text = raw_bytes.decode("utf-8")
    except UnicodeDecodeError:
        raise HTTPException(status_code=400, detail="Could not decode file. Please upload a UTF-8 encoded .txt or .md file.")

    slides = _run_generation(doc_text, slide_instructions, ai_model, slide_count, content_length, sarvam_key, gemini_key)

    # Translate if needed
    if target_language != "en-IN":
        if not sarvam_key:
            raise HTTPException(status_code=422, detail="sarvam_key is required for translation.")
        for s in slides:
            s["heading"] = translate_content(s["heading"], target_language, sarvam_key)
            for item in s["content"]:
                item["text"] = translate_content(item["text"], target_language, sarvam_key)

    return SlidesResponse(slides=slides, slide_count=len(slides), ai_model=ai_model)


@app.post("/generate/pptx", tags=["Generate"])
async def generate_pptx(
    doc_file: UploadFile = File(..., description="Documentation file (.txt or .md)"),
    template_file: UploadFile = File(..., description="PowerPoint template (.pptx)"),
    slide_count: int = Form(5, ge=2, le=10, description="Number of slides to generate"),
    content_length: str = Form("Medium", description="'Brief', 'Medium', or 'Detailed'"),
    ai_model: str = Form("Sarvam AI", description="'Sarvam AI' or 'Google Gemini'"),
    sarvam_key: str = Form("", description="Sarvam API key (required for Sarvam AI or translation)"),
    gemini_key: str = Form("", description="Google Gemini API key (required for Google Gemini)"),
    slide_instructions: str = Form("", description="Optional per-slide instructions"),
    target_language: str = Form("en-IN", description="Target language code (e.g. 'hi-IN')"),
):
    """
    Generate a complete PowerPoint file from uploaded documentation and a `.pptx` template.
    
    Returns the finished `.pptx` file as a direct download.
    """
    if content_length not in CONTENT_LENGTH_MAP:
        raise HTTPException(status_code=422, detail=f"content_length must be one of: {list(CONTENT_LENGTH_MAP.keys())}")
    if target_language not in SUPPORTED_LANGUAGES:
        raise HTTPException(status_code=422, detail=f"Unsupported language code. Use /languages to see valid options.")

    raw_doc = await doc_file.read()
    raw_template = await template_file.read()

    try:
        doc_text = raw_doc.decode("utf-8")
    except UnicodeDecodeError:
        raise HTTPException(status_code=400, detail="Could not decode documentation file. Please use UTF-8 encoding.")

    slides = _run_generation(doc_text, slide_instructions, ai_model, slide_count, content_length, sarvam_key, gemini_key)

    # Translate if needed
    if target_language != "en-IN":
        if not sarvam_key:
            raise HTTPException(status_code=422, detail="sarvam_key is required for translation.")
        for s in slides:
            s["heading"] = translate_content(s["heading"], target_language, sarvam_key)
            for item in s["content"]:
                item["text"] = translate_content(item["text"], target_language, sarvam_key)

    pptx_bytes = _build_pptx(slides, raw_template)

    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=presentation.pptx"},
    )
